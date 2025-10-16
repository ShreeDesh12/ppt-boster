"""
PowerPoint presentation generator service
"""
import os
import logging
from typing import List, Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from app.models import SlideContent, SlideLayout, ThemeConfig, Citation
from app.config import get_settings

logger = logging.getLogger(__name__)
settings = get_settings()


class PresentationGenerator:
    """Generates PowerPoint presentations"""

    def __init__(self, theme: Optional[ThemeConfig] = None):
        logger.info("Initializing PresentationGenerator")
        self.theme = theme or ThemeConfig()
        logger.debug(f"Theme configured: primary_color={self.theme.primary_color}, font={self.theme.font_name}")
        self.prs = None

    def create_presentation(
        self,
        slides: List[SlideContent],
        citations: Optional[List[Citation]] = None,
        aspect_ratio: str = "16:9"
    ) -> Presentation:
        """
        Create a PowerPoint presentation from slide content

        Args:
            slides: List of slide content
            citations: Optional list of citations
            aspect_ratio: Presentation aspect ratio (16:9 or 4:3)

        Returns:
            Presentation object
        """
        logger.info(f"Creating presentation with {len(slides)} slides, aspect_ratio={aspect_ratio}")

        self.prs = Presentation()

        # Set slide dimensions based on aspect ratio
        if aspect_ratio == "16:9":
            self.prs.slide_width = Inches(10)
            self.prs.slide_height = Inches(5.625)
            logger.debug("Set aspect ratio to 16:9 (10 x 5.625 inches)")
        else:  # 4:3
            self.prs.slide_width = Inches(10)
            self.prs.slide_height = Inches(7.5)
            logger.debug("Set aspect ratio to 4:3 (10 x 7.5 inches)")

        # Generate each slide
        for idx, slide_content in enumerate(slides, 1):
            logger.info(f"Adding slide {idx}/{len(slides)}: {slide_content.layout.value} - {slide_content.title[:50]}")
            self._add_slide(slide_content)

        # Add citations slide if provided
        if citations:
            logger.info(f"Adding citations slide with {len(citations)} citations")
            self._add_citations_slide(citations)

        logger.info("Presentation creation completed")
        return self.prs

    def save_presentation(self, file_path: str) -> str:
        """
        Save presentation to file

        Args:
            file_path: Path to save the presentation

        Returns:
            Absolute path to saved file
        """
        logger.info(f"Saving presentation to: {file_path}")

        directory = os.path.dirname(file_path)
        if directory:
            os.makedirs(directory, exist_ok=True)
            logger.debug(f"Ensured directory exists: {directory}")

        self.prs.save(file_path)
        abs_path = os.path.abspath(file_path)
        logger.info(f"Presentation saved successfully: {abs_path}")
        return abs_path

    def _add_slide(self, content: SlideContent):
        """Add a slide based on layout type"""
        logger.debug(f"_add_slide called with layout={content.layout.value}")

        if content.layout == SlideLayout.TITLE:
            self._add_title_slide(content)
        elif content.layout == SlideLayout.BULLET_POINTS:
            self._add_bullet_points_slide(content)
        elif content.layout == SlideLayout.TWO_COLUMN:
            self._add_two_column_slide(content)
        elif content.layout == SlideLayout.CONTENT_WITH_IMAGE:
            self._add_content_with_image_slide(content)

        logger.debug(f"Slide added successfully: {content.layout.value}")

    def _add_title_slide(self, content: SlideContent):
        """Add a title slide"""
        logger.debug(f"_add_title_slide: title='{content.title[:50]}', content='{content.content[:50] if content.content else 'None'}'")
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank layout

        # Add title
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(1)

        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        title_frame.text = content.title

        title_paragraph = title_frame.paragraphs[0]
        title_paragraph.alignment = PP_ALIGN.CENTER
        title_paragraph.font.size = Pt(self.theme.font_size_title)
        title_paragraph.font.name = self.theme.font_name
        title_paragraph.font.bold = True
        title_paragraph.font.color.rgb = self._hex_to_rgb(self.theme.primary_color)

        # Add subtitle if content exists
        if content.content:
            subtitle_top = Inches(3.2)
            subtitle_box = slide.shapes.add_textbox(left, subtitle_top, width, Inches(0.8))
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = content.content

            subtitle_paragraph = subtitle_frame.paragraphs[0]
            subtitle_paragraph.alignment = PP_ALIGN.CENTER
            subtitle_paragraph.font.size = Pt(self.theme.font_size_body)
            subtitle_paragraph.font.name = self.theme.font_name

    def _add_bullet_points_slide(self, content: SlideContent):
        """Add a bullet points slide"""
        logger.debug(f"_add_bullet_points_slide: title='{content.title[:50]}', {len(content.bullet_points) if content.bullet_points else 0} bullets")
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = content.title

        title_paragraph = title_frame.paragraphs[0]
        title_paragraph.font.size = Pt(self.theme.font_size_title - 8)
        title_paragraph.font.name = self.theme.font_name
        title_paragraph.font.bold = True
        title_paragraph.font.color.rgb = self._hex_to_rgb(self.theme.primary_color)

        # Add bullet points
        left = Inches(1)
        top = Inches(1.8)
        width = Inches(8)
        height = Inches(3.5)

        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        text_frame.word_wrap = True

        if content.bullet_points:
            for i, bullet in enumerate(content.bullet_points):
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()

                p.text = bullet
                p.level = 0
                p.font.size = Pt(self.theme.font_size_body)
                p.font.name = self.theme.font_name
                p.space_before = Pt(12)

    def _add_two_column_slide(self, content: SlideContent):
        """Add a two-column layout slide"""
        logger.debug(f"_add_two_column_slide: title='{content.title[:50]}'")
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = content.title

        title_paragraph = title_frame.paragraphs[0]
        title_paragraph.font.size = Pt(self.theme.font_size_title - 8)
        title_paragraph.font.name = self.theme.font_name
        title_paragraph.font.bold = True
        title_paragraph.font.color.rgb = self._hex_to_rgb(self.theme.primary_color)

        # Left column
        left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(4.5), Inches(3.5))
        left_frame = left_box.text_frame
        left_frame.word_wrap = True
        left_frame.text = content.left_column or ""

        left_paragraph = left_frame.paragraphs[0]
        left_paragraph.font.size = Pt(self.theme.font_size_body)
        left_paragraph.font.name = self.theme.font_name

        # Right column
        right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.8), Inches(4.5), Inches(3.5))
        right_frame = right_box.text_frame
        right_frame.word_wrap = True
        right_frame.text = content.right_column or ""

        right_paragraph = right_frame.paragraphs[0]
        right_paragraph.font.size = Pt(self.theme.font_size_body)
        right_paragraph.font.name = self.theme.font_name

    def _add_content_with_image_slide(self, content: SlideContent):
        """Add a content with image placeholder slide"""
        logger.debug(f"_add_content_with_image_slide: title='{content.title[:50]}', has_image_desc={bool(content.image_description)}")
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = content.title

        title_paragraph = title_frame.paragraphs[0]
        title_paragraph.font.size = Pt(self.theme.font_size_title - 8)
        title_paragraph.font.name = self.theme.font_name
        title_paragraph.font.bold = True
        title_paragraph.font.color.rgb = self._hex_to_rgb(self.theme.primary_color)

        # Add content text
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(5), Inches(3.5))
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        content_frame.text = content.content or ""

        content_paragraph = content_frame.paragraphs[0]
        content_paragraph.font.size = Pt(self.theme.font_size_body)
        content_paragraph.font.name = self.theme.font_name

        # Add image placeholder
        left = Inches(6)
        top = Inches(1.8)
        width = Inches(3.5)
        height = Inches(3)

        shape = slide.shapes.add_shape(
            1,  # Rectangle
            left, top, width, height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(220, 220, 220)

        # Add image description text
        if content.image_description:
            text_frame = shape.text_frame
            text_frame.text = f"[Image: {content.image_description}]"
            text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            text_frame.paragraphs[0].font.size = Pt(12)
            text_frame.paragraphs[0].font.italic = True

    def _add_citations_slide(self, citations: List[Citation]):
        """Add a citations/references slide"""
        logger.debug(f"_add_citations_slide: adding {len(citations)} citations")
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = "References"

        title_paragraph = title_frame.paragraphs[0]
        title_paragraph.font.size = Pt(self.theme.font_size_title - 8)
        title_paragraph.font.name = self.theme.font_name
        title_paragraph.font.bold = True
        title_paragraph.font.color.rgb = self._hex_to_rgb(self.theme.primary_color)

        # Add citations
        text_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9), Inches(3.5))
        text_frame = text_box.text_frame
        text_frame.word_wrap = True

        for i, citation in enumerate(citations):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()

            citation_text = citation.source
            if citation.title:
                citation_text += f" - {citation.title}"
            if citation.date:
                citation_text += f" ({citation.date})"

            p.text = citation_text
            p.font.size = Pt(self.theme.font_size_body - 2)
            p.font.name = self.theme.font_name
            p.space_before = Pt(8)

    def _hex_to_rgb(self, hex_color: str) -> RGBColor:
        """Convert hex color to RGB"""
        logger.debug(f"_hex_to_rgb: converting {hex_color} to RGB")
        hex_color = hex_color.lstrip('#')
        rgb = RGBColor(
            int(hex_color[0:2], 16),
            int(hex_color[2:4], 16),
            int(hex_color[4:6], 16)
        )
        logger.debug(f"RGB result: R={rgb}, G={rgb}, B={rgb}")
        return rgb
